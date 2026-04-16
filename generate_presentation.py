# generate_presentation.py — NGSTability Presentation with LaTeX Equations
# Requirements: pip install python-pptx matplotlib
# Run: python generate_presentation.py

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, tempfile, shutil

# --- Equation rendering via matplotlib ---
try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False
    print("WARNING: matplotlib not found. Equations will render as plain text.")
    print("  Install with: pip install matplotlib")

EQ_DIR = tempfile.mkdtemp(prefix='ngstab_eq_') if HAS_MATPLOTLIB else None
_eq_counter = [0]

def render_equation(latex_str, fontsize=18, dpi=200, color='#333333'):
    """Render a LaTeX string to a transparent PNG. Returns path or None."""
    if not HAS_MATPLOTLIB:
        return None
    _eq_counter[0] += 1
    fname = os.path.join(EQ_DIR, f'eq_{_eq_counter[0]:03d}.png')
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.set_axis_off()
    ax.text(0, 0, f'${latex_str}$', fontsize=fontsize, color=color,
            ha='left', va='baseline', transform=ax.transAxes)
    fig.savefig(fname, dpi=dpi, transparent=True, bbox_inches='tight', pad_inches=0.05)
    plt.close(fig)
    return fname

def add_eq_image(slide, latex_str, left_in, top_in, height_in=0.5, fontsize=20, color='#333333'):
    """Render LaTeX and place as image on slide. Falls back to plain text."""
    img = render_equation(latex_str, fontsize=fontsize, color=color)
    if img:
        slide.shapes.add_picture(img, Inches(left_in), Inches(top_in), height=Inches(height_in))
    else:
        # Fallback: show raw-ish text
        fallback = latex_str.replace(r'\cdot', '\u00b7').replace(r'\frac', '').replace(r'\quad', '  ')
        for tag in [r'\left', r'\right', r'\mathrm', r'\,', r'\!', r'\;']:
            fallback = fallback.replace(tag, '')
        fallback = fallback.replace('{', '').replace('}', '')
        tx = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(10), Inches(height_in))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = fallback
        p.font.size = Pt(15); p.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
MED_BLUE  = RGBColor(0x00, 0x70, 0xC0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GRAY      = RGBColor(0x40, 0x40, 0x40)
GREEN     = RGBColor(0x00, 0x80, 0x00)
RED       = RGBColor(0xCC, 0x00, 0x00)
TBL_HDR   = RGBColor(0x00, 0x50, 0x90)
TBL_ROW1  = RGBColor(0xE8, 0xF0, 0xFA)
TBL_ROW2  = RGBColor(0xFF, 0xFF, 0xFF)

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.LEFT
    tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.5))
    tf2 = tx2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(18); p2.font.color.rgb = RGBColor(0xBB,0xCC,0xDD); p2.alignment = PP_ALIGN.LEFT
    return slide

def add_content_slide(title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0xFF,0xFF,0xFF))
    # title bar
    tb = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))  # rectangle
    tb.fill.solid(); tb.fill.fore_color.rgb = DARK_BLUE; tb.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = tx.text_frame; p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE
    # bullets
    bx = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    btf = bx.text_frame; btf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = b; p.font.size = Pt(16); p.font.color.rgb = GRAY; p.space_after = Pt(8)
        p.level = 0
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = btf.add_paragraph()
                sp.text = sb; sp.font.size = Pt(14); sp.font.color.rgb = RGBColor(0x66,0x66,0x66)
                sp.space_after = Pt(4); sp.level = 1
    return slide

def add_table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    tb = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    tb.fill.solid(); tb.fill.fore_color.rgb = DARK_BLUE; tb.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = tx.text_frame; p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = Inches(11.5) if not col_widths else sum(Inches(w) for w in col_widths)
    tbl_h = Inches(0.4) * n_rows
    top = Inches(1.5)
    left = Inches(0.9)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_w, tbl_h)
    table = shape.table

    if col_widths:
        for ci, w in enumerate(col_widths):
            table.columns[ci].width = Inches(w)

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(13); para.font.bold = True; para.font.color.rgb = WHITE; para.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HDR

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri+1, ci)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12); para.font.color.rgb = BLACK
            cell.fill.solid()
            cell.fill.fore_color.rgb = TBL_ROW1 if ri % 2 == 0 else TBL_ROW2

    return slide

def make_content_slide(title):
    """Create a white slide with dark-blue title bar. Returns the slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))
    tb = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    tb.fill.solid(); tb.fill.fore_color.rgb = DARK_BLUE; tb.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = tx.text_frame; p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE
    return slide

def add_bullets(slide, bullets, left=0.8, top=1.4, width=11.5, height=5.5, size=16, color=GRAY):
    """Place bullet text at a specific position on a slide."""
    bx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    btf = bx.text_frame; btf.word_wrap = True
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = b; p.font.size = Pt(size); p.font.color.rgb = color; p.space_after = Pt(6)

# ============================================================
# TITLE SLIDE
# ============================================================
add_title_slide(
    "NGSTability \u2014 Next Gen PCM Buck Converter\nStability Analysis & PID Optimization",
    "Digital Control Loop Modeling  |  RTL-Accurate Coefficients  |  Global Optimization  |  AMS Correlation"
)

# ============================================================
# MOTIVATION SLIDE
# ============================================================
slide = make_content_slide("Why NGSTability?  \u2014  Motivation")
add_bullets(slide, [
    "New architecture: moved from voltage-mode analog compensation to PCM buck with digital PID",
    "   \u2192  Entirely different control loop \u2014 no legacy model or design methodology to reuse",
    "",
    "The gaps we faced:",
    "   \u2022 No accurate model to predict which K1/K2/K3 PID sets produce a stable loop",
    "   \u2022 PID tuning was done manually by trial-and-error \u2014 both in lab and in AMS simulation",
    "   \u2022 Even after finding a stable set, we had no visibility into actual PM/GM margins",
    "   \u2022 Pre-silicon: each AMS trial takes hours \u2014 searching for a stable set is painfully slow",
    "",
    "NGSTability approach:",
    "   1.  Model the full closed-loop system in MATLAB (Laplace domain)",
    "        \u2192  Bode plot \u2192 extract GM, PM, crossover frequency in milliseconds",
    "   2.  Validate: if MATLAB predicts stable, it should correlate with lab and AMS measurements",
    "   3.  Once the model is trusted, run global optimization over all possible PID sets",
    "        \u2192  Find the best K1/K2/K3 under stability and hardware constraints automatically",
    "   4.  Replaces days of manual AMS iteration with minutes of automated search",
], left=0.8, top=1.3, height=5.8, size=15)

# ============================================================
# SLIDE 1.1 \u2014 PCM Feedback Loop Architecture
# ============================================================
add_table_slide(
    "1. Small-Signal Modeling \u2014 PCM Feedback Loop Architecture",
    ["Block", "Transfer Function", "Key Parameters"],
    [
        ["PCM Plant", "Gvc = Hdc \u00b7 Fp \u00b7 Fh", "L=0.47\u00b5H, C=22\u00b5F, Resr=3m\u03a9"],
        ["Digital PID", "Hpid = Kp + Ki/s + Kd\u00b7s/(Tf\u00b7s+1)", "After >>6 shift"],
        ["ADC delay", "Pad\u00e9(5ns, order 2)", "4-bit, 80mV window"],
        ["DAC filter", "1 / (1 + s/\u03c9dac)", "500 kHz pole, 12-bit"],
        ["Blanking + Comp delay", "Pad\u00e9(15ns + 8ns)", ""],
        ["Feedback", "H = \u03b2 = 0.5", "Resistive divider"],
    ],
    col_widths=[2.5, 4.5, 4.5]
)

# SLIDE 1.2 — PCM Plant Model (with LaTeX equations)
slide = make_content_slide("1. Small-Signal Modeling \u2014 PCM Plant: Control-to-Output")
add_eq_image(slide,
    r'G_{vc}(s) = H_{dc} \cdot F_p(s) \cdot F_h(s)',
    1.0, 1.35, height_in=0.45, fontsize=22)
add_eq_image(slide,
    r'H_{dc} = \frac{R}{R_i} \cdot \frac{1}{1 + \frac{R \, T_s \, (m_c D^\prime - 0.5)}{L}}',
    1.0, 2.0, height_in=0.75, fontsize=20)
add_eq_image(slide,
    r'F_p(s) = \frac{1 + s \, C \, R_{esr}}{1 + s / \omega_p}',
    1.0, 2.95, height_in=0.6, fontsize=20)
add_eq_image(slide,
    r'F_h(s) = \frac{1}{1 + \dfrac{s}{\omega_n Q_p} + \dfrac{s^2}{\omega_n^2}}',
    1.0, 3.75, height_in=0.8, fontsize=20)
add_bullets(slide, [
    "Slope compensation:  mc = 1 + se/sn,   se = 0.7 MV/s",
    "Plant gain varies with load current (1 mA \u2013 3 A)  \u2192  must sweep full range",
], left=0.8, top=4.8, height=1.5, size=15)

# ============================================================
# SLIDE 2.1 \u2014 Closed-Form Components
# ============================================================
slide = make_content_slide("2. Closed-Form Loop Components \u2014 Transfer Functions")

add_bullets(slide, ["Analog reference compensator (Type-II):"], left=0.8, top=1.3, height=0.4, size=16)
add_eq_image(slide,
    r'H_c(s) = \frac{\omega_{p1}}{s} \cdot \frac{1 + s/\omega_{z,comp}}{1 + s/\omega_{p2}}',
    1.2, 1.7, height_in=0.55, fontsize=18)

add_bullets(slide, ["Digital PID compensator (Laplace domain):"], left=0.8, top=2.35, height=0.4, size=16)
add_eq_image(slide,
    r'H_{pid}(s) = K_p + \frac{K_i}{s} + \frac{K_d \cdot s}{T_f \cdot s + 1}',
    1.2, 2.7, height_in=0.55, fontsize=18)

add_bullets(slide, ["Complete digital compensator chain:"], left=0.8, top=3.35, height=0.4, size=16)
add_eq_image(slide,
    r'H_{dig}(s) = G_{ADC} \cdot H_{adc}(s) \cdot H_{pid}(s) \cdot G_{digital} \cdot G_{DAC} \cdot H_{dac}(s) \cdot H_{blank}(s) \cdot e^{-sT_d}',
    1.2, 3.7, height_in=0.35, fontsize=16)

add_bullets(slide, ["Loop gains and closed-loop:"], left=0.8, top=4.2, height=0.4, size=16)
add_eq_image(slide,
    r'L_c(s) = H_c(s) \cdot G_{vc}(s) \cdot H \qquad \textrm{(analog reference)}',
    1.2, 4.55, height_in=0.3, fontsize=16)
add_eq_image(slide,
    r'L_{loop}(s) = H_{dig}(s) \cdot G_{vc}(s) \cdot H \qquad \textrm{(digital loop gain)}',
    1.2, 4.95, height_in=0.3, fontsize=16)
add_eq_image(slide,
    r'T_{loop}(s) = \frac{L_{loop}(s)}{1 + L_{loop}(s)} \qquad \textrm{(closed-loop)}',
    1.2, 5.35, height_in=0.5, fontsize=16)

add_bullets(slide, [
    "Delays modeled as 2nd-order Pad\u00e9     |     G_digital = 2\u207b\u2076 = 1/64     |     ADC/DAC gains = 1",
], left=0.8, top=6.0, height=0.5, size=13)

# SLIDE 2.2 \u2014 Phase Budget
add_table_slide(
    "2. Phase Budget at Crossover (fc = 600 kHz)",
    ["Block", "Phase Contribution", "Notes"],
    [
        ["Plant \u00d7 H", "~ \u2212XX\u00b0", "Load-dependent"],
        ["ADC delay (5 ns)", "~ \u2212X\u00b0", "Pad\u00e9 model"],
        ["Digital PID", "~ +XX\u00b0", "Phase boost from zeros"],
        ["DAC filter (500 kHz)", "~ \u2212X\u00b0", "Single pole"],
        ["Blanking delay (15 ns)", "~ \u2212X\u00b0", "Pad\u00e9 model"],
        ["Comp delay (8 ns)", "~ \u2212X\u00b0", "Pad\u00e9 model"],
        ["TOTAL", "~ \u2212XXX\u00b0", ""],
        ["Phase Margin", "PM = Total + 180\u00b0", "Target \u2265 35\u00b0"],
    ],
    col_widths=[3.5, 3.5, 4.5]
)

# ============================================================
# SLIDE 3.1 \u2014 RTL Coefficient Framework
# ============================================================
slide = make_content_slide("3. RTL-Accurate Coefficient Mapping")
add_bullets(slide, [
    "Three equivalent representations of the same controller:",
    "   Lab Position [P, I, D]  \u2194  RTL [K1, K2, K3]  \u2194  Effective [Kp, Ki, Kd]",
], left=0.8, top=1.3, height=0.7, size=16)
add_bullets(slide, ["Lab \u2192 RTL:"], left=0.8, top=2.1, height=0.3, size=15)
add_eq_image(slide,
    r'K_1 = P + I + D \qquad K_2 = -P - 2D \qquad K_3 = D',
    1.5, 2.4, height_in=0.3, fontsize=18)
add_bullets(slide, ["RTL \u2192 PID:"], left=0.8, top=2.8, height=0.3, size=15)
add_eq_image(slide,
    r'K_p = -K_2 - 2K_3 \qquad K_i = \frac{K_1 + K_2 + K_3}{T_{ctrl}} \qquad K_d = K_3 \cdot T_{ctrl}',
    1.5, 3.1, height_in=0.5, fontsize=18)
add_bullets(slide, ["Effective (post >>6 shift):"], left=0.8, top=3.7, height=0.3, size=15)
add_eq_image(slide,
    r'K_{p,eff} = K_p \cdot G_{dig} \qquad K_{i,eff} = K_i \cdot G_{dig} \qquad K_{d,eff} = K_d \cdot G_{dig}',
    1.5, 4.0, height_in=0.3, fontsize=18)
add_bullets(slide, ["RTL PID recursion (velocity form):"], left=0.8, top=4.4, height=0.3, size=15)
add_eq_image(slide,
    r'u[n] = u[n\!-\!1] + K_1 \cdot e[n] + K_2 \cdot e[n\!-\!1] + K_3 \cdot e[n\!-\!2]',
    1.5, 4.7, height_in=0.3, fontsize=18)
add_eq_image(slide,
    r'\mathrm{output} = (u \gg 6) + 2048',
    1.5, 5.1, height_in=0.3, fontsize=18)
add_bullets(slide, [
    "All conversions are invertible \u2014 any representation maps to any other",
], left=0.8, top=5.5, height=0.5, size=14)

# SLIDE 3.2 \u2014 Quantization Constraints
add_table_slide(
    "3. Hardware Constraints on PID Design Space",
    ["Coefficient", "Bit Width", "Range", "Constraint"],
    [
        ["K1", "20-bit signed", "\u00b1524,287", "|K1| \u2264 2^19 \u2212 1"],
        ["K2", "20-bit signed", "\u00b1524,287", "|K2| \u2264 2^19 \u2212 1"],
        ["K3", "11-bit signed", "\u00b11,023", "|K3| \u2264 2^10 \u2212 1"],
        ["PID accumulator", "Pre-clamp", "\u2212131,072 to +131,008", "RTL MIN/MAX_VALUE"],
        ["PID output", "12-bit unsigned", "0 \u2013 4,095", "After >>6 + 2048 offset"],
        ["Ki resolution", "~4\u00d710\u2076 LSB", "Coarse steps", "Integer I_cmd quantization"],
    ],
    col_widths=[2.8, 2.5, 3.2, 3.0]
)

# ============================================================
# SLIDE 4.1 \u2014 Optimization Architecture
# ============================================================
slide = make_content_slide("4. Global PID Optimization \u2014 Architecture")
add_bullets(slide, [
    "Goal: Find [Kp, Ki, Kd] maximizing stability + transient performance across 9 loads",
    "",
    "Search space (log\u2081\u2080 scaled):",
], left=0.8, top=1.3, height=1.0, size=16)
add_eq_image(slide,
    r'K_p \in [0.016,\; 10^4] \qquad K_i \in [4 \times 10^6,\; 5 \times 10^8] \qquad K_d \in [K_{d,min},\; K_{d,max}]',
    1.2, 2.35, height_in=0.4, fontsize=17)
add_bullets(slide, [
    "Load sweep:  Iload = [1mA, 100mA, 300mA, 500mA, 700mA, 900mA, 1.1A, 2A, 3A]",
    "",
    "Two-stage optimizer:",
    "   1. PSO (Particle Swarm) \u2014 40 particles \u00d7 80 iterations \u2192 global exploration",
    "   2. Nelder-Mead hybrid refinement \u2192 local polishing of best PSO result",
    "",
    "Hard constraints:",
    "   \u2022 Closed-loop stable at ALL load points",
    "   \u2022 PM \u2265 35\u00b0,  GM \u2265 6 dB at every load",
    "   \u2022 K1/K2/K3 within RTL bit-width (200-point penalty per overflow)",
], left=0.8, top=2.85, height=4.0, size=15)

# SLIDE 4.2 \u2014 Cost Function
slide = make_content_slide("4. Multi-Metric Cost Function")
add_eq_image(slide,
    r'\mathrm{cost} = -\,\mathrm{score} + \mathrm{overflow\_penalty}',
    1.0, 1.3, height_in=0.35, fontsize=20)
add_bullets(slide, ["Score = stability margins + load-step transient metrics:"], left=0.8, top=1.75, height=0.3, size=15)
add_eq_image(slide,
    r'\mathrm{score} = \underbrace{2 \cdot PM_{min} + 1 \cdot GM_{min} + 0.15 \cdot PM_{avg} + 0.1 \cdot GM_{avg} - 1 \cdot \overline{\left|\log_{10}\!\tfrac{f_c}{f_{c,t}}\right|}}_{\textrm{stability margins}}',
    1.0, 2.1, height_in=0.75, fontsize=15)
add_eq_image(slide,
    r'+ \underbrace{12 \cdot \min\!\left(\frac{\mathrm{slew}}{0.002},\; 10\right) - 12 \cdot \frac{t_{settle}}{5\,\mu s} - 12 \cdot \frac{\Delta V_{us}}{20\,\mathrm{mV}}}_{\textrm{load-step transient}}',
    1.0, 3.0, height_in=0.75, fontsize=15)
add_bullets(slide, [
    "",
    "All terms contribute to a single score \u2014 optimizer maximizes the total",
    "",
    "Margin caps:  PM capped at 85\u00b0,  GM capped at 40 dB",
    "Overflow penalty:  +200 per K1/K2/K3 exceeding RTL bit-width",
    "",
    "Convergence visualization: cost vs iteration with PSO \u2192 fminsearch transition marker",
], left=0.8, top=3.9, height=3.0, size=15)

# ============================================================
# SLIDE 5.1 \u2014 Digital vs Analog Comparison
# ============================================================
add_table_slide(
    "5. Optimized Digital Loop vs Ideal Analog Compensator",
    ["Metric", "Analog (Ideal Hc)", "Digital (Optimized)", "Notes"],
    [
        ["Compensator type", "Type-II continuous", "Velocity PID (K1/K2/K3)", ""],
        ["Phase Margin", "XX\u00b0", "XX\u00b0", "Fill from simulation"],
        ["Gain Margin", "XX dB", "XX dB", "Fill from simulation"],
        ["Crossover freq", "600 kHz", "~XXX kHz", "Target = fs/10"],
        ["Phase loss sources", "None", "ADC/DAC/blanking delays", "10\u201320\u00b0 typical loss"],
        ["Quantization", "Infinite resolution", "Integer K1/K2/K3", "Finite bit-width"],
        ["Load robustness", "Single operating pt", "9 loads simultaneously", "1mA to 3A"],
    ],
    col_widths=[2.5, 3.0, 3.0, 3.0]
)

# SLIDE 5.2 \u2014 3D Stability Map
add_content_slide(
    "5. 3D Stability Design Space Visualization",
    [
        "Axes: Kd (X),  Ki (Y),  Kp (Z)  \u2014 log-scaled",
        "",
        "GREEN = stable at ALL 9 load points (PM \u2265 35\u00b0, GM \u2265 6 dB)",
        "RED = unstable or insufficient margins at any load point",
        "BLUE STAR = optimizer's selected solution",
        "",
        "Key insights:",
        "   \u2022 The stable region is a bounded volume \u2014 not all PID combinations work",
        "   \u2022 Optimizer solution sits inside the stable region, near the best-margin boundary",
        "   \u2022 Boundaries shift with load current \u2014 multi-load check is essential",
        "",
        "2D projections (Kd-Ki, Kd-Kp, Ki-Kp) show the stable envelope for each pair",
        "   \u2192 Green if stable for ANY value in the collapsed axis",
        "",
        "Grid: 25 \u00d7 25 \u00d7 15 = 9,375 points evaluated",
    ]
)

# ============================================================
# SLIDE 6.1 \u2014 Nonlinear Predictor
# ============================================================
add_table_slide(
    "6. RTL-Accurate Cycle-by-Cycle Nonlinear Predictor",
    ["Feature", "Implementation", "Details"],
    [
        ["ADC", "4-bit quantization", "0.36\u20130.44V window, 16 levels, 5mV/LSB"],
        ["PID accumulator", "Velocity form + pre-clamp", "K1\u00b7e[n] + K2\u00b7e[n\u22121] + K3\u00b7e[n\u22122]"],
        ["Shift + offset", ">>6 then +2048", "Unsigned 12-bit output"],
        ["DAC", "12-bit code \u2192 analog", "NMOS Vth cap at 1.6V (not 1.8V)"],
        ["DAC filter", "500 kHz analog pole", "Discrete IIR at Tctrl"],
        ["Plant model", "State-space @ Tctrl", "Control path + load disturbance (parallel)"],
        ["Load disturbance", "Gload = \u2212Zout_OL", "Passive output impedance only"],
        ["Load step", "Configurable profile", "Delay, rise time, amplitude, width"],
    ],
    col_widths=[2.5, 3.5, 5.5]
)

# SLIDE 6.2 \u2014 Validation Pipeline
slide = make_content_slide("6. Validation Pipeline: Predictor \u2192 AMS \u2192 Silicon")
add_bullets(slide, [
    "Tiered validation flow:",
    "   Optimizer  \u2192  Nonlinear Predictor  \u2192  AMS Simulation  \u2192  Silicon",
    "   (fast, 1000s of evals)    (quant/saturation)      (full circuit)       (final proof)",
], left=0.8, top=1.3, height=1.0, size=15)
add_bullets(slide, ["Load disturbance model:"], left=0.8, top=2.5, height=0.3, size=15)
add_eq_image(slide,
    r'G_{load} = -Z_{out,OL} = -\frac{Z_C \cdot R}{Z_C + R}',
    1.2, 2.8, height_in=0.6, fontsize=20)
add_bullets(slide, [
    "Metrics compared at each stage:",
    "   \u2022 Vout undershoot:   target < 20 mV",
    "   \u2022 Recovery slew:     target > 0.002 V/\u00b5s",
    "   \u2022 Settling time:     target < 5 \u00b5s",
    "   \u2022 Steady-state Vout: 0.800 V",
    "",
    "This tiered approach saves weeks of AMS iteration time",
], left=0.8, top=3.55, height=3.3, size=15)

# ============================================================
# SUMMARY SLIDE
# ============================================================
add_content_slide(
    "Summary \u2014 NGSTability Methodology",
    [
        "1.  Small-signal Laplace-domain model of the PCM feedback loop",
        "     \u2192 Closed-form plant, digital blocks, delays \u2014 all analytically derived",
        "",
        "2.  RTL-accurate coefficient mapping (K1/K2/K3 \u2194 Lab P/I/D \u2194 Kp/Ki/Kd)",
        "     \u2192 Exact quantization and bit-width constraints from RTL",
        "",
        "3.  Global optimization (PSO + Nelder-Mead) across full load range",
        "     \u2192 Cost function: PM + GM + crossover + load-step transient metrics",
        "",
        "4.  3D stability map visualization of the feasible PID design space",
        "",
        "5.  Cycle-by-cycle nonlinear predictor for AMS correlation",
        "     \u2192 RTL-accurate ADC/DAC quantization, PID clamping, load disturbance",
        "",
        "6.  Validated pipeline: Linear model \u2192 Predictor \u2192 AMS \u2192 Silicon",
        "     \u2192 Each tier catches progressively finer effects while managing compute cost",
    ]
)

# ============================================================
# SAVE
# ============================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "NGSTability_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")

# Clean up temporary equation images
if EQ_DIR and os.path.isdir(EQ_DIR):
    shutil.rmtree(EQ_DIR, ignore_errors=True)
    print(f"Cleaned up {_eq_counter[0]} equation images.")
